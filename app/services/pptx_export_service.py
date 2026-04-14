"""Structured PPTX export helpers for lesson content."""

from __future__ import annotations

from dataclasses import dataclass, field
from io import BytesIO
from pathlib import Path
import re
from typing import Any, Callable, Dict, Iterable, List, Optional

from bs4 import BeautifulSoup

from app.schemas.api_models import LessonImageResource
from src.observability.logger import get_logger

logger = get_logger(__name__)

_PPTX_IMPORTS: Optional[Dict[str, Any]] = None

_NOTE_PREFIXES = ("讲解提示：", "教师提示：", "教师讲解：", "备注：")
_PRACTICE_KEYWORDS = ("练习", "检测", "巩固", "任务", "活动", "思考", "探究")
_SUMMARY_KEYWORDS = ("小结", "总结", "回顾", "板书", "归纳")
_VISUAL_KEYWORDS = ("导入", "情境", "案例", "实验", "图", "示意", "现象", "比较")


@dataclass
class LessonSlide:
    title: str
    layout: str = "standard"
    bullets: List[str] = field(default_factory=list)
    paragraphs: List[str] = field(default_factory=list)
    speaker_notes: List[str] = field(default_factory=list)
    image_sources: List[str] = field(default_factory=list)
    accent_text: Optional[str] = None


@dataclass
class LessonDeck:
    title: str
    slides: List[LessonSlide] = field(default_factory=list)


def _get_pptx_imports() -> Dict[str, Any]:
    global _PPTX_IMPORTS
    if _PPTX_IMPORTS is not None:
        return _PPTX_IMPORTS

    try:
        from pptx import Presentation
        from pptx.dml.color import RGBColor
        from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
        from pptx.util import Inches, Pt
    except ImportError as e:
        raise RuntimeError("PPTX 导出依赖缺失，请安装 python-pptx 后重试") from e

    _PPTX_IMPORTS = {
        "Presentation": Presentation,
        "RGBColor": RGBColor,
        "MSO_AUTO_SHAPE_TYPE": MSO_AUTO_SHAPE_TYPE,
        "Inches": Inches,
        "Pt": Pt,
    }
    return _PPTX_IMPORTS


def _clean_text(text: str) -> str:
    return re.sub(r"\s+", " ", str(text or "")).strip()


def _normalize_slide_title(text: str) -> str:
    cleaned = _clean_text(text)
    cleaned = re.sub(r"^幻灯片\d+\s*[｜|]\s*", "", cleaned)
    return cleaned or "课件内容"


def _collect_list_items(tag: Any) -> List[str]:
    items: List[str] = []
    for item in tag.find_all("li", recursive=False):
        text = _clean_text(item.get_text(" ", strip=True))
        if text:
            items.append(text)
    return items


def _iter_image_sources(image_resources: Iterable[LessonImageResource | Dict[str, Any]]) -> List[str]:
    urls: List[str] = []
    for item in image_resources or []:
        if isinstance(item, LessonImageResource):
            url = _clean_text(item.url)
        elif isinstance(item, dict):
            url = _clean_text(item.get("url") or "")
        else:
            url = _clean_text(getattr(item, "url", "") or "")
        if url and url not in urls:
            urls.append(url)
    return urls


def _extract_note(text: str) -> Optional[str]:
    normalized = _clean_text(text)
    for prefix in _NOTE_PREFIXES:
        if normalized.startswith(prefix):
            return normalized[len(prefix):].strip() or None
    return None


def _infer_layout(title: str, bullets: List[str], paragraphs: List[str], image_sources: List[str]) -> str:
    text = " ".join([title, *bullets, *paragraphs]).lower()
    if "封面" in title:
        return "cover"
    if any(keyword in text for keyword in _SUMMARY_KEYWORDS):
        return "summary"
    if any(keyword in text for keyword in _PRACTICE_KEYWORDS):
        return "practice"
    if image_sources or any(keyword in text for keyword in _VISUAL_KEYWORDS):
        return "two_column"
    return "standard"


def _chunk_slide(slide: LessonSlide) -> List[LessonSlide]:
    content_items = [{"kind": "bullet", "text": item} for item in slide.bullets] + [
        {"kind": "paragraph", "text": item} for item in slide.paragraphs
    ]
    if not content_items:
        slide.layout = _infer_layout(slide.title, slide.bullets, slide.paragraphs, slide.image_sources)
        return [slide]

    limit = 5 if slide.layout in {"practice", "summary"} else 6
    chunks = [content_items[i:i + limit] for i in range(0, len(content_items), limit)]
    chunked: List[LessonSlide] = []
    for index, chunk in enumerate(chunks):
        chunked.append(
            LessonSlide(
                title=slide.title if index == 0 else f"{slide.title}（续）",
                layout=slide.layout,
                bullets=[item["text"] for item in chunk if item["kind"] == "bullet"],
                paragraphs=[item["text"] for item in chunk if item["kind"] == "paragraph"],
                speaker_notes=list(slide.speaker_notes if index == 0 else []),
                image_sources=list(slide.image_sources if index == 0 else []),
                accent_text=slide.accent_text if index == 0 else None,
            )
        )
    return chunked


def _allocate_fallback_images(slides: List[LessonSlide], available_image_sources: List[str]) -> None:
    used = {src for slide in slides for src in slide.image_sources}
    remaining = [src for src in available_image_sources if src not in used]
    if not remaining:
        return

    preferred_indexes: List[int] = []
    for idx, slide in enumerate(slides):
        if slide.layout == "cover" or slide.image_sources:
            continue
        text = " ".join([slide.title, *slide.bullets, *slide.paragraphs])
        if slide.layout == "two_column" or any(keyword in text for keyword in _VISUAL_KEYWORDS):
            preferred_indexes.append(idx)

    if not preferred_indexes:
        preferred_indexes = [
            idx for idx, slide in enumerate(slides)
            if slide.layout in {"standard", "two_column"} and not slide.image_sources
        ]

    for idx in preferred_indexes:
        if not remaining:
            break
        slides[idx].image_sources.append(remaining.pop(0))
        if slides[idx].layout == "standard":
            slides[idx].layout = "two_column"


def build_lesson_deck(
    *,
    title: str,
    content_html: str,
    image_resources: Optional[List[LessonImageResource | Dict[str, Any]]] = None,
) -> LessonDeck:
    soup = BeautifulSoup(f"<div>{content_html}</div>", "html.parser")
    root = soup.div
    if root is None:
        return LessonDeck(title=_clean_text(title), slides=[LessonSlide(title=_clean_text(title), layout="cover")])

    deck_title = _clean_text(title)
    current_slide: Optional[LessonSlide] = None
    slides: List[LessonSlide] = []

    def commit_current() -> None:
        nonlocal current_slide
        if current_slide is None:
            return
        current_slide.layout = _infer_layout(
            current_slide.title,
            current_slide.bullets,
            current_slide.paragraphs,
            current_slide.image_sources,
        )
        slides.extend(_chunk_slide(current_slide))
        current_slide = None

    for node in root.children:
        if not getattr(node, "name", None):
            continue

        name = str(node.name).lower()
        text = _clean_text(node.get_text(" ", strip=True))

        if name == "h1":
            if text:
                deck_title = text
            continue

        if name == "h2":
            commit_current()
            current_slide = LessonSlide(title=_normalize_slide_title(text))
            continue

        if current_slide is None:
            current_slide = LessonSlide(title="课件内容")

        note = _extract_note(text) if text else None
        if note:
            current_slide.speaker_notes.append(note)
            continue

        if name in {"ul", "ol"}:
            current_slide.bullets.extend(_collect_list_items(node))
            continue

        if name == "h3":
            if text:
                if current_slide.accent_text is None:
                    current_slide.accent_text = text
                else:
                    current_slide.bullets.append(text)
            continue

        if name == "blockquote":
            if text:
                current_slide.speaker_notes.append(text)
            continue

        if name == "p":
            images = node.find_all("img")
            if images:
                for image in images:
                    src = _clean_text(image.get("src") or "")
                    if src and src not in current_slide.image_sources:
                        current_slide.image_sources.append(src)
                remaining_text = _clean_text(node.get_text(" ", strip=True))
                if remaining_text and not remaining_text.startswith("配图"):
                    current_slide.paragraphs.append(remaining_text)
                continue

            if text:
                current_slide.paragraphs.append(text)
            continue

        if name == "img":
            src = _clean_text(node.get("src") or "")
            if src and src not in current_slide.image_sources:
                current_slide.image_sources.append(src)

    commit_current()

    if not slides:
        slides = [LessonSlide(title=_clean_text(title) or "课件封面", layout="cover")]

    if slides[0].layout != "cover":
        slides.insert(0, LessonSlide(title=deck_title, layout="cover"))
    elif not slides[0].title:
        slides[0].title = deck_title

    _allocate_fallback_images(slides, _iter_image_sources(image_resources or []))
    return LessonDeck(title=deck_title, slides=slides)


def _add_textbox(
    slide: Any,
    left: Any,
    top: Any,
    width: Any,
    height: Any,
    text: str,
    *,
    font_size: int,
    bold: bool = False,
    color: tuple[int, int, int] = (34, 34, 34),
) -> Any:
    imports = _get_pptx_imports()
    textbox = slide.shapes.add_textbox(left, top, width, height)
    frame = textbox.text_frame
    frame.clear()
    paragraph = frame.paragraphs[0]
    paragraph.text = text
    run = paragraph.runs[0]
    run.font.name = "Microsoft YaHei"
    run.font.size = imports["Pt"](font_size)
    run.font.bold = bold
    run.font.color.rgb = imports["RGBColor"](*color)
    return textbox


def _add_bullet_box(slide: Any, items: List[str], *, left: Any, top: Any, width: Any, height: Any, font_size: int = 18) -> None:
    imports = _get_pptx_imports()
    textbox = slide.shapes.add_textbox(left, top, width, height)
    frame = textbox.text_frame
    frame.clear()
    frame.word_wrap = True
    for index, item in enumerate(items):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.text = item
        paragraph.level = 0
        paragraph.bullet = True
        run = paragraph.runs[0]
        run.font.name = "Microsoft YaHei"
        run.font.size = imports["Pt"](font_size)
        run.font.color.rgb = imports["RGBColor"](35, 38, 45)


def _set_speaker_notes(slide: Any, notes: List[str]) -> None:
    clean_notes = [_clean_text(note) for note in notes if _clean_text(note)]
    if not clean_notes:
        return
    notes_frame = slide.notes_slide.notes_text_frame
    notes_frame.text = "\n".join(clean_notes)


def _add_footer(slide: Any, text: str = "Modular RAG Lesson PPT") -> None:
    imports = _get_pptx_imports()
    footer = slide.shapes.add_textbox(
        imports["Inches"](0.75),
        imports["Inches"](6.9),
        imports["Inches"](11.2),
        imports["Inches"](0.3),
    )
    frame = footer.text_frame
    frame.clear()
    p = frame.paragraphs[0]
    p.text = text
    run = p.runs[0]
    run.font.name = "Microsoft YaHei"
    run.font.size = imports["Pt"](9)
    run.font.color.rgb = imports["RGBColor"](125, 125, 125)


def _add_header_bar(slide: Any, title: str, *, accent: tuple[int, int, int]) -> None:
    imports = _get_pptx_imports()
    header = slide.shapes.add_shape(
        imports["MSO_AUTO_SHAPE_TYPE"].RECTANGLE,
        0,
        0,
        imports["Inches"](13.333),
        imports["Inches"](0.62),
    )
    header.fill.solid()
    header.fill.fore_color.rgb = imports["RGBColor"](*accent)
    header.line.fill.background()
    _add_textbox(
        slide,
        imports["Inches"](0.7),
        imports["Inches"](0.15),
        imports["Inches"](11.3),
        imports["Inches"](0.4),
        title,
        font_size=20,
        bold=True,
        color=(255, 255, 255),
    )


def _add_picture_if_available(
    slide: Any,
    *,
    image_sources: List[str],
    resolve_image_path: Callable[[str], Optional[Path]],
    left: Any,
    top: Any,
    width: Any,
    height: Any,
) -> bool:
    for src in image_sources:
        image_path = resolve_image_path(src)
        if image_path is None or not image_path.exists():
            continue
        try:
            slide.shapes.add_picture(str(image_path), left, top, width=width, height=height)
            return True
        except Exception as exc:
            logger.warning("lesson_pptx.add_picture_failed src=%s path=%s error=%r", src, image_path, exc)
    return False


def _render_cover_slide(presentation: Any, lesson_slide: LessonSlide, deck_title: str) -> None:
    imports = _get_pptx_imports()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])

    bg = slide.shapes.add_shape(
        imports["MSO_AUTO_SHAPE_TYPE"].RECTANGLE,
        0,
        0,
        presentation.slide_width,
        presentation.slide_height,
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = imports["RGBColor"](243, 247, 252)
    bg.line.fill.background()

    stripe = slide.shapes.add_shape(
        imports["MSO_AUTO_SHAPE_TYPE"].RECTANGLE,
        imports["Inches"](0.75),
        imports["Inches"](0.9),
        imports["Inches"](0.2),
        imports["Inches"](2.0),
    )
    stripe.fill.solid()
    stripe.fill.fore_color.rgb = imports["RGBColor"](29, 95, 173)
    stripe.line.fill.background()

    _add_textbox(
        slide,
        imports["Inches"](1.15),
        imports["Inches"](1.05),
        imports["Inches"](10.5),
        imports["Inches"](1.4),
        lesson_slide.title or deck_title,
        font_size=26,
        bold=True,
        color=(27, 43, 71),
    )
    subtitle_lines = lesson_slide.bullets[:3] or lesson_slide.paragraphs[:2]
    if subtitle_lines:
        _add_bullet_box(
            slide,
            subtitle_lines,
            left=imports["Inches"](1.2),
            top=imports["Inches"](2.45),
            width=imports["Inches"](6.4),
            height=imports["Inches"](2.3),
            font_size=18,
        )
    _set_speaker_notes(slide, lesson_slide.speaker_notes)
    _add_footer(slide, "封面页")


def _render_two_column_slide(presentation: Any, lesson_slide: LessonSlide, resolve_image_path: Callable[[str], Optional[Path]]) -> None:
    imports = _get_pptx_imports()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    _add_header_bar(slide, lesson_slide.title, accent=(29, 95, 173))

    if lesson_slide.accent_text:
        _add_textbox(
            slide,
            imports["Inches"](0.8),
            imports["Inches"](0.95),
            imports["Inches"](5.9),
            imports["Inches"](0.45),
            lesson_slide.accent_text,
            font_size=16,
            bold=True,
            color=(29, 95, 173),
        )

    content_items = lesson_slide.bullets or lesson_slide.paragraphs
    _add_bullet_box(
        slide,
        content_items[:6],
        left=imports["Inches"](0.8),
        top=imports["Inches"](1.35),
        width=imports["Inches"](5.8),
        height=imports["Inches"](5.1),
        font_size=18,
    )

    card = slide.shapes.add_shape(
        imports["MSO_AUTO_SHAPE_TYPE"].ROUNDED_RECTANGLE,
        imports["Inches"](7.1),
        imports["Inches"](1.15),
        imports["Inches"](5.35),
        imports["Inches"](4.1),
    )
    card.fill.solid()
    card.fill.fore_color.rgb = imports["RGBColor"](247, 249, 252)
    card.line.color.rgb = imports["RGBColor"](213, 222, 234)

    has_picture = _add_picture_if_available(
        slide,
        image_sources=lesson_slide.image_sources,
        resolve_image_path=resolve_image_path,
        left=imports["Inches"](7.25),
        top=imports["Inches"](1.3),
        width=imports["Inches"](5.05),
        height=imports["Inches"](3.7),
    )
    if not has_picture:
        _add_textbox(
            slide,
            imports["Inches"](7.45),
            imports["Inches"](2.25),
            imports["Inches"](4.6),
            imports["Inches"](0.9),
            "此页适合放置示意图、案例图或流程图",
            font_size=16,
            bold=True,
            color=(98, 112, 135),
        )

    note_hint = lesson_slide.paragraphs[0] if lesson_slide.paragraphs and lesson_slide.paragraphs != content_items else None
    if note_hint:
        _add_textbox(
            slide,
            imports["Inches"](7.25),
            imports["Inches"](5.45),
            imports["Inches"](5.0),
            imports["Inches"](0.8),
            note_hint,
            font_size=12,
            color=(90, 98, 112),
        )
    _set_speaker_notes(slide, lesson_slide.speaker_notes)
    _add_footer(slide)


def _render_practice_slide(presentation: Any, lesson_slide: LessonSlide) -> None:
    imports = _get_pptx_imports()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    _add_header_bar(slide, lesson_slide.title, accent=(187, 101, 25))

    for idx, item in enumerate((lesson_slide.bullets or lesson_slide.paragraphs)[:5], start=1):
        top = 0.95 + (idx - 1) * 1.1
        box = slide.shapes.add_shape(
            imports["MSO_AUTO_SHAPE_TYPE"].ROUNDED_RECTANGLE,
            imports["Inches"](0.85),
            imports["Inches"](top),
            imports["Inches"](11.6),
            imports["Inches"](0.78),
        )
        box.fill.solid()
        box.fill.fore_color.rgb = imports["RGBColor"](252, 247, 240)
        box.line.color.rgb = imports["RGBColor"](234, 208, 183)
        _add_textbox(
            slide,
            imports["Inches"](1.05),
            imports["Inches"](top + 0.12),
            imports["Inches"](11.0),
            imports["Inches"](0.5),
            f"{idx}. {item}",
            font_size=17,
            color=(68, 54, 33),
        )

    if lesson_slide.accent_text:
        _add_textbox(
            slide,
            imports["Inches"](0.95),
            imports["Inches"](6.3),
            imports["Inches"](10.8),
            imports["Inches"](0.45),
            lesson_slide.accent_text,
            font_size=13,
            bold=True,
            color=(160, 88, 20),
        )
    _set_speaker_notes(slide, lesson_slide.speaker_notes)
    _add_footer(slide, "练习与互动")


def _render_summary_slide(presentation: Any, lesson_slide: LessonSlide) -> None:
    imports = _get_pptx_imports()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    _add_header_bar(slide, lesson_slide.title, accent=(65, 133, 83))

    panel = slide.shapes.add_shape(
        imports["MSO_AUTO_SHAPE_TYPE"].ROUNDED_RECTANGLE,
        imports["Inches"](0.95),
        imports["Inches"](1.15),
        imports["Inches"](11.3),
        imports["Inches"](4.8),
    )
    panel.fill.solid()
    panel.fill.fore_color.rgb = imports["RGBColor"](244, 250, 245)
    panel.line.color.rgb = imports["RGBColor"](194, 219, 198)

    _add_bullet_box(
        slide,
        (lesson_slide.bullets or lesson_slide.paragraphs)[:6],
        left=imports["Inches"](1.25),
        top=imports["Inches"](1.5),
        width=imports["Inches"](10.2),
        height=imports["Inches"](3.9),
        font_size=20,
    )

    if lesson_slide.accent_text:
        _add_textbox(
            slide,
            imports["Inches"](1.15),
            imports["Inches"](6.15),
            imports["Inches"](10.5),
            imports["Inches"](0.4),
            lesson_slide.accent_text,
            font_size=13,
            bold=True,
            color=(65, 133, 83),
        )
    _set_speaker_notes(slide, lesson_slide.speaker_notes)
    _add_footer(slide, "总结与回顾")


def _render_standard_slide(presentation: Any, lesson_slide: LessonSlide) -> None:
    imports = _get_pptx_imports()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    _add_header_bar(slide, lesson_slide.title, accent=(74, 82, 140))

    if lesson_slide.accent_text:
        _add_textbox(
            slide,
            imports["Inches"](0.85),
            imports["Inches"](0.95),
            imports["Inches"](11.2),
            imports["Inches"](0.45),
            lesson_slide.accent_text,
            font_size=15,
            bold=True,
            color=(74, 82, 140),
        )

    items = lesson_slide.bullets or lesson_slide.paragraphs
    _add_bullet_box(
        slide,
        items[:6],
        left=imports["Inches"](0.95),
        top=imports["Inches"](1.35),
        width=imports["Inches"](11.1),
        height=imports["Inches"](4.9),
        font_size=19,
    )
    _set_speaker_notes(slide, lesson_slide.speaker_notes)
    _add_footer(slide)


def _render_lesson_slide(
    presentation: Any,
    lesson_slide: LessonSlide,
    *,
    deck_title: str,
    resolve_image_path: Callable[[str], Optional[Path]],
) -> None:
    layout = lesson_slide.layout
    if layout == "cover":
        _render_cover_slide(presentation, lesson_slide, deck_title)
        return
    if layout == "two_column":
        _render_two_column_slide(presentation, lesson_slide, resolve_image_path)
        return
    if layout == "practice":
        _render_practice_slide(presentation, lesson_slide)
        return
    if layout == "summary":
        _render_summary_slide(presentation, lesson_slide)
        return
    _render_standard_slide(presentation, lesson_slide)


def build_lesson_pptx_bytes(
    *,
    title: str,
    content_html: str,
    resolve_image_path: Callable[[str], Optional[Path]],
    image_resources: Optional[List[LessonImageResource | Dict[str, Any]]] = None,
) -> bytes:
    imports = _get_pptx_imports()
    deck = build_lesson_deck(title=title, content_html=content_html, image_resources=image_resources or [])

    presentation = imports["Presentation"]()
    presentation.slide_width = imports["Inches"](13.333)
    presentation.slide_height = imports["Inches"](7.5)

    for lesson_slide in deck.slides:
        _render_lesson_slide(
            presentation,
            lesson_slide,
            deck_title=deck.title,
            resolve_image_path=resolve_image_path,
        )

    buffer = BytesIO()
    presentation.save(buffer)
    return buffer.getvalue()
