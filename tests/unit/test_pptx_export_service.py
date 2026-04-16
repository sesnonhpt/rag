from __future__ import annotations

from io import BytesIO
from pathlib import Path

from pptx import Presentation

from app.schemas.ppt_models import LessonDeckModel, LessonSlideModel, SlideElementModel
from app.services.pptx_export_service import (
    build_lesson_deck,
    build_lesson_pptx_bytes,
    build_lesson_pptx_bytes_from_deck,
)


def test_build_lesson_pptx_bytes_returns_zip_payload(tmp_path: Path):
    image_path = tmp_path / "demo.png"
    image_path.write_bytes(
        bytes.fromhex(
            "89504E470D0A1A0A0000000D4948445200000001000000010802000000907753DE"
            "0000000C49444154789C63606060000000040001F61738550000000049454E44AE426082"
        )
    )

    pptx_bytes = build_lesson_pptx_bytes(
        title="测试课件",
        content_html=(
            "<h1>《牛顿第三定律》PPT课件</h1>"
            "<h2>幻灯片1｜学习目标</h2>"
            "<ul><li>理解作用力与反作用力</li><li>能区分平衡力</li></ul>"
            "<p><img src='/demo.png' alt='配图1' /></p>"
            "<h2>幻灯片2｜课堂小结</h2>"
            "<p>大小相等，方向相反，作用在两个物体上。</p>"
        ),
        resolve_image_path=lambda src: image_path if src == "/demo.png" else None,
    )

    assert pptx_bytes.startswith(b"PK")
    assert len(pptx_bytes) > 1000


def test_build_lesson_deck_assigns_layouts_and_fallback_images():
    deck = build_lesson_deck(
        title="测试课件",
        content_html=(
            "<h1>《牛顿第三定律》PPT课件</h1>"
            "<h2>幻灯片1｜封面</h2>"
            "<ul><li>高中物理</li><li>新授课</li></ul>"
            "<h2>幻灯片2｜案例分析</h2>"
            "<ul><li>掰手腕情境</li><li>滑冰现象</li></ul>"
            "<p>讲解提示：先让学生判断施力物体和受力物体。</p>"
            "<h2>幻灯片3｜课堂练习</h2>"
            "<ul><li>判断是否属于作用力与反作用力</li></ul>"
        ),
        image_resources=[{"url": "/lesson-plan-image/demo-image"}],
    )

    assert deck.slides[0].layout == "cover"
    assert deck.slides[1].layout == "two_column"
    assert deck.slides[1].image_sources == ["/lesson-plan-image/demo-image"]
    assert "先让学生判断" in deck.slides[1].speaker_notes[0]
    assert deck.slides[2].layout == "practice"


def test_build_lesson_pptx_bytes_writes_speaker_notes(tmp_path: Path):
    image_path = tmp_path / "demo.png"
    image_path.write_bytes(
        bytes.fromhex(
            "89504E470D0A1A0A0000000D4948445200000001000000010802000000907753DE"
            "0000000C49444154789C63606060000000040001F61738550000000049454E44AE426082"
        )
    )

    pptx_bytes = build_lesson_pptx_bytes(
        title="测试课件",
        content_html=(
            "<h1>《牛顿第三定律》PPT课件</h1>"
            "<h2>幻灯片1｜封面</h2>"
            "<ul><li>高中物理</li></ul>"
            "<h2>幻灯片2｜案例分析</h2>"
            "<ul><li>掰手腕情境</li></ul>"
            "<p>讲解提示：提醒学生比较两个物体上的力。</p>"
        ),
        resolve_image_path=lambda src: image_path if src == "/lesson-plan-image/demo-image" else None,
        image_resources=[{"url": "/lesson-plan-image/demo-image"}],
    )

    prs = Presentation(BytesIO(pptx_bytes))
    assert len(prs.slides) >= 2
    assert "提醒学生比较两个物体上的力" in prs.slides[1].notes_slide.notes_text_frame.text


def test_build_lesson_pptx_bytes_from_deck_uses_elements():
    deck = LessonDeckModel(
        deck_id="deck_elements",
        title="元素课件",
        slides=[
            LessonSlideModel(
                id="slide_1",
                order=0,
                title="元素页",
                layout="standard",
                speaker_notes=["元素备注"],
                elements=[
                    SlideElementModel(
                        id="bg",
                        type="shape",
                        x=0,
                        y=0,
                        w=1000,
                        h=46,
                        fill_color="#1d5fad",
                    ),
                    SlideElementModel(
                        id="title",
                        type="text",
                        x=56,
                        y=10,
                        w=820,
                        h=28,
                        text="元素标题",
                        font_size=22,
                        bold=True,
                        text_color="#ffffff",
                    ),
                ],
            )
        ],
    )

    pptx_bytes = build_lesson_pptx_bytes_from_deck(deck=deck, resolve_image_path=lambda src: None)
    prs = Presentation(BytesIO(pptx_bytes))
    assert len(prs.slides) == 1
    assert "元素备注" in prs.slides[0].notes_slide.notes_text_frame.text
